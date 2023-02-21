# 使用from...import从pptx模块中导入Presentation
from pptx import Presentation

# 将.pptx文件路径赋值给变量path
path = "C:/Users/ChuanZhou/Desktop/python 学习运用/应用工具/pptx/statistics.pptx"

# 读取path并赋值给变量pptxFile
pptxFile = Presentation(path)

# for循环遍历pptxFile中的.slides属性，赋值给slide
for slide in pptxFile.slides:

    # for循环遍历slide中.shapes属性，赋值给变量shape
    for shape in slide.shapes:
        # 判断形状中是否有文本框
        if shape.has_text_frame == True:
            # 读取形状中的文本框，并赋值给变量textFrame
            textFrame = shape.text_frame
            
            # for循环遍历文本框内的所有段落
            # 赋值给变量paragraph
            for paragraph in textFrame.paragraphs:
                for run in paragraph.runs:
                    # 读取样式块中的文本内容，并赋值给变量texts
                    texts = run.text
                    # print()输出texts
                    print(texts)
                